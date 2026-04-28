#!/usr/bin/env python3
"""
Surgical PPTX builder:
1. Run all fixes via python-pptx to get modified slide XMLs
2. Start from the ORIGINAL zip (preserves everything byte-for-byte)
3. Replace ONLY the 5 modified slide files (4, 5, 6, 13, 19)
4. Save — zero risk of corrupting slides 7-30

This eliminates the python-pptx save() rewrite problem.
"""
import sys, io, os, zipfile
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

BASE = r'c:\Users\amont\Desktop\PYG\investor-ready\01-Pitch-Deck'
TMP  = r'C:\Users\amont\Documents'

DECKS = [
    ('PYGLARA_Pitch_Deck_EN.pptx', 'PYGLARA_Pitch_Deck_EN_v3.pptx', 'en'),
    ('PYGLARA_Pitch_Deck.pptx',    'PYGLARA_Pitch_Deck_ES_v3.pptx', 'es'),
]

# ── Import all fix functions from the existing scripts ──────────────────────
sys.path.insert(0, r'C:\Users\amont\Desktop\PYG')
from fix_deck_v3      import fix_slide4, fix_slide5, fix_slide6
from fix_deck_s13_s19 import fix_slide12, fix_slide13, fix_slide14, fix_slide19


def get_modified_slides(src_path, lang):
    """
    Run all fixes on a copy, then extract ONLY the 5 modified slide XMLs.
    Returns dict: {slide_zip_path: xml_bytes}
    """
    from pptx import Presentation

    prs = Presentation(src_path)

    fix_slide4(prs.slides[3], lang)
    fix_slide5(prs.slides[4], lang)
    fix_slide6(prs.slides[5], lang)
    fix_slide12(prs.slides[11], lang)
    fix_slide13(prs.slides[12], lang)
    fix_slide14(prs.slides[13], lang)
    fix_slide19(prs.slides[18], lang)

    # Save to a temp buffer
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    # Extract only the 5 modified slide XML files from the temp output
    modified = {}
    with zipfile.ZipFile(buf, 'r') as z:
        # We need to find which ZIP entries correspond to slides 4,5,6,13,19
        # python-pptx may have remapped them — use sldIdLst order
        import xml.etree.ElementTree as ET
        prs_xml = ET.fromstring(z.read('ppt/presentation.xml'))
        rels_xml = ET.fromstring(z.read('ppt/_rels/presentation.xml.rels'))

        NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        NS_R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
        SLD  = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide'

        rid2target = {r.get('Id'): r.get('Target')
                      for r in rels_xml if r.get('Type') == SLD}

        sld_lst = prs_xml.find(f'.//{{{NS_P}}}sldIdLst')
        ordered = [rid2target[el.get(f'{{{NS_R}}}id')]
                   for el in sld_lst]

        # Position 3,4,5,12,18 (0-indexed) = slides 4,5,6,13,19
        for pos in [3, 4, 5, 11, 12, 13, 18]:
            rel_target = ordered[pos]   # e.g. 'slides/slide4.xml'
            zip_path   = f'ppt/{rel_target}'
            xml_bytes  = z.read(zip_path)
            # Map to the ORIGINAL zip path (we'll figure out which is which below)
            modified[pos] = (zip_path, xml_bytes)
            print(f'  Captured pos {pos+1} -> {zip_path} ({len(xml_bytes)} bytes)')

    return modified


def build_surgical_pptx(orig_path, dst_path, lang):
    print(f'\n=== Building {os.path.basename(dst_path)} ===')

    # Step 1: Get modified slide XMLs from python-pptx
    print('Step 1: Generating modified slides...')
    modified = get_modified_slides(orig_path, lang)

    # Step 2: Read original zip
    print('Step 2: Reading original zip...')
    with open(orig_path, 'rb') as fh:
        orig_bytes = fh.read()

    # Step 3: Determine original slide paths for positions 4,5,6,13,19
    print('Step 3: Mapping positions to original zip paths...')
    with zipfile.ZipFile(io.BytesIO(orig_bytes), 'r') as orig_z:
        import xml.etree.ElementTree as ET
        prs_xml  = ET.fromstring(orig_z.read('ppt/presentation.xml'))
        rels_xml = ET.fromstring(orig_z.read('ppt/_rels/presentation.xml.rels'))

        NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        NS_R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
        SLD  = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide'

        rid2target = {r.get('Id'): r.get('Target')
                      for r in rels_xml if r.get('Type') == SLD}
        sld_lst = prs_xml.find(f'.//{{{NS_P}}}sldIdLst')
        ordered = [rid2target[el.get(f'{{{NS_R}}}id')] for el in sld_lst]

        # Build replacement map: orig_zip_path -> new_xml_bytes
        replacements = {}
        for pos, (_, xml_bytes) in modified.items():
            orig_zip_path = 'ppt/' + ordered[pos]
            replacements[orig_zip_path] = xml_bytes
            print(f'  Replace {orig_zip_path} ({len(xml_bytes)} bytes)')

        # Step 4: Rebuild ZIP from original, only replacing those 5 files
        print('Step 4: Rebuilding zip with surgical replacements...')
        buf = io.BytesIO()
        with zipfile.ZipFile(io.BytesIO(orig_bytes), 'r') as src_z, \
             zipfile.ZipFile(buf, 'w') as dst_z:

            for info in src_z.infolist():
                if info.filename in replacements:
                    # Write replaced content
                    dst_z.writestr(info, replacements[info.filename])
                    print(f'  [REPLACED] {info.filename}')
                else:
                    # Copy original bytes exactly
                    dst_z.writestr(info, src_z.read(info.filename))

    # Step 5: Write to destination
    with open(dst_path, 'wb') as fh:
        fh.write(buf.getvalue())

    size_kb = os.path.getsize(dst_path) // 1024
    print(f'Step 5: Saved {os.path.basename(dst_path)} ({size_kb} KB)')

    # Step 6: Verify
    from pptx import Presentation
    p = Presentation(dst_path)
    slides = len(p.slides)
    slide_map = {}
    for i, slide in enumerate(p.slides):
        texts = [t.strip() for s in slide.shapes if s.has_text_frame
                 for t in [pp.text for pp in s.text_frame.paragraphs] if t.strip()]
        slide_map[i+1] = texts[0][:45] if texts else '(empty)'

    print(f'\n  Verification — {slides} slides:')
    for num, title in slide_map.items():
        print(f'  Slide {num:2d}: {title}')


for orig_name, dst_name, lang in DECKS:
    orig_path = os.path.join(BASE, orig_name)
    dst_path  = os.path.join(BASE, dst_name)
    build_surgical_pptx(orig_path, dst_path, lang)

print('\n\nAll done.')
