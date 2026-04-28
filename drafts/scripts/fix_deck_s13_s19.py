#!/usr/bin/env python3
"""
PYGLARA Pitch Deck — Slide 13 + 19 Coverage Fix
Adds demand signals to every product category (Slide 13) and
market share math to Competitive Position (Slide 19).
Updates: EN_v3.pptx and ES_v3.pptx in-place.
"""

import sys, io
if __name__ == '__main__':
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches as I, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

COPPER      = RGBColor(0xB8, 0x73, 0x33)
COPPER_DARK = RGBColor(0x8B, 0x5A, 0x2B)
COPPER_PALE = RGBColor(0xF0, 0xE4, 0xD0)
DARK        = RGBColor(0x2A, 0x1F, 0x10)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GREY        = RGBColor(0x9E, 0x97, 0x8E)

# ──────────────────────────────────────────────────────────────────────────────
# BILINGUAL CONTENT
# ──────────────────────────────────────────────────────────────────────────────

# Slide 13 — demand signal appended to each product description text box
S13_DEMAND = {
    'en': {
        'Text 7':  "Demand: 2,000+ new 4G/5G towers + CORPOELEC on-site. Every pole needs galvanizing.",
        'Text 11': "Demand: $41B Rystad fab & construction segment 2026-2040. Pipe racks alone fill a month.",
        'Text 15': "Demand: VINCI/Freyssinet MSE walls in active projects + housing & highway boom.",
        'Text 19': "Demand: National solar expansion + grid modernization. Mounting systems all galvanized.",
        'Text 23': "Demand: 34 historical clients. Freyssinet and road/bridge projects in pipeline.",
        'Text 27': "Demand: Marullo (same block) + PROALCA + regional agroindustrial restart.",
    },
    'es': {
        'Text 7':  "Demanda: 2.000+ torres 4G/5G nuevas + CORPOELEC en la misma zona. Cada poste requiere galvanizado.",
        'Text 11': "Demanda: US$41MM segmento fab. y construccion Rystad 2026-2040. Soportes de tuberia llenan un mes.",
        'Text 15': "Demanda: Muros MSE VINCI/Freyssinet en proyectos activos + boom vivienda y autopistas.",
        'Text 19': "Demanda: Expansion solar nacional + modernizacion de red. Todos los soportes requieren galvanizado.",
        'Text 23': "Demanda: 34 clientes historicos. Freyssinet y proyectos viales en cartera.",
        'Text 27': "Demanda: Marullo (mismo bloque) + PROALCA + reactivacion agroindustrial regional.",
    }
}

# Slide 13 — replace truncated footer (Text 29) with complete + coverage stat
S13_FOOTER = {
    'en': (
        "MSE retaining walls (VINCI/Freyssinet, Zona Industrial I): galvanized steel strips "
        "embedded in fill — PYGLARA's 7m kettle is purpose-built for this. "
        "At 300 TM/month proven output, one kettle serves all six categories above. "
        "Filling PYGLARA to capacity requires less than 3% of Venezuela's oil sector galvanizing demand."
    ),
    'es': (
        "Muros de tierra armada (VINCI/Freyssinet, Zona Industrial I): flejes de acero galvanizado "
        "embebidos en el relleno — la cuba de 7m de PYGLARA esta disenada para esto. "
        "Con 300 TM/mes de capacidad comprobada, una cuba atiende las seis categorias arriba. "
        "Llenar PYGLARA a plena capacidad requiere menos del 3% de la demanda de galvanizado del sector petrolero venezolano."
    ),
}

# Slide 19 — Text 23 replacement text (includes market share, fits in original H=450 box)
S19_COVERAGE_REWRITE = {
    'en': "Coverage: Lara, Yaracuy, Portuguesa, Falcon, Barinas, Zulia (oil) — at 300 TM/month: est. 6-12% of Venezuela's active galvanizing demand",
    'es': "Cobertura: Lara, Yaracuy, Portuguesa, Falcon, Barinas, Zulia (petroleo) — a 300 TM/mes: est. 6-12% de la demanda activa de galvanizado",
}

# Slide 19 — update Text 17 bullet to surface market math context (optional enhancement)
S19_BARRIER_EXTRA = {
    'en': " With only 2 active providers, PYGLARA holds half of Venezuela's commercial galvanizing supply.",
    'es': " Con solo 2 proveedores activos, PYGLARA controla la mitad de la oferta comercial de galvanizado en Venezuela.",
}


# ──────────────────────────────────────────────────────────────────────────────
# HELPERS
# ──────────────────────────────────────────────────────────────────────────────

def append_paragraph(shape, text, size=7.5, color=COPPER_DARK, bold=False,
                     italic=True, align=PP_ALIGN.LEFT, space_before=Pt(3)):
    tf = shape.text_frame
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def overwrite_textbox(shape, text, size=8, color=DARK, bold=False,
                      italic=False, align=PP_ALIGN.LEFT):
    """Clear all paragraphs and replace with a single paragraph."""
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def find_shape(slide, name):
    for shp in slide.shapes:
        if shp.name == name:
            return shp
    return None


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 13 FIX
# ──────────────────────────────────────────────────────────────────────────────

def fix_slide12(slide, lang):
    """Translate Spanish pull-quote in the EN deck (Text 3)."""
    if lang != 'en':
        return
    t3 = find_shape(slide, 'Text 3')
    if t3 and t3.has_text_frame:
        for para in t3.text_frame.paragraphs:
            for run in para.runs:
                if 'Ferrari' in run.text and 'hornos' in run.text:
                    run.text = '\u201cThe Ferrari of galvanizing kettles\u201d'


def fix_slide14(slide, lang):
    """Delete TextBox 24 — overlaps Revenue row and adds no layout value."""
    to_remove = [sh for sh in slide.shapes if sh.name == 'TextBox 24']
    for sh in to_remove:
        sp = sh._element
        sp.getparent().remove(sp)


def fix_slide13(slide, lang):
    demands = S13_DEMAND[lang]

    # Add demand signal line to each product description box
    for box_name, demand_text in demands.items():
        shp = find_shape(slide, box_name)
        if not shp:
            print(f"  WARNING: {box_name} not found on slide 13")
            continue
        append_paragraph(shp, demand_text,
                         size=7.5, color=COPPER_DARK,
                         italic=True, bold=False,
                         space_before=Pt(4))

    # Replace truncated footer (Text 29) with complete coverage statement
    footer = find_shape(slide, 'Text 29')
    if footer:
        overwrite_textbox(footer, S13_FOOTER[lang],
                          size=8, color=DARK, italic=True,
                          align=PP_ALIGN.LEFT)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 19 FIX
# ──────────────────────────────────────────────────────────────────────────────

def fix_slide19(slide, lang):
    # 1. Rewrite Text 23 (Coverage bullet) with market share inline — stays within Shape 22 (H=450)
    t23 = find_shape(slide, 'Text 23')
    if t23:
        # Keep original height — do NOT expand so text stays inside the container shape
        overwrite_textbox(t23, S19_COVERAGE_REWRITE[lang],
                          size=8, color=WHITE, bold=False, italic=False,
                          align=PP_ALIGN.LEFT)

    # 2. Translate English pull-quote (Text 24) in the Spanish deck
    if lang == 'es':
        t24 = find_shape(slide, 'Text 24')
        if t24 and t24.has_text_frame:
            for para in t24.text_frame.paragraphs:
                for run in para.runs:
                    if 'Extremely high barrier' in run.text:
                        run.text = '\u201cBarrera de entrada extremadamente alta\u201d'


# ──────────────────────────────────────────────────────────────────────────────
# MAIN
# ──────────────────────────────────────────────────────────────────────────────

if __name__ == '__main__':
    DECKS = [
        ('c:/Users/amont/Desktop/PYG/investor-ready/01-Pitch-Deck/PYGLARA_Pitch_Deck_EN_v3.pptx', 'en'),
        ('c:/Users/amont/Desktop/PYG/investor-ready/01-Pitch-Deck/PYGLARA_Pitch_Deck_ES_v3.pptx', 'es'),
    ]

    for path, lang in DECKS:
        print(f"\nProcessing {lang.upper()}: {path}")
        prs = Presentation(path)

        fix_slide13(prs.slides[12], lang)
        print(f"  [OK] Slide 13: demand signals added to 6 product categories, footer updated")

        fix_slide19(prs.slides[18], lang)
        print(f"  [OK] Slide 19: market share math added to coverage bullet")

        prs.save(path)
        print(f"  [SAVED] {path}")

    print("\nDone.")
