#!/usr/bin/env python3
"""
PYGLARA Pitch Deck v3 - Automated Fix Script
Fixes: Slide 4 (color revert, timeline rebuild, remove Regime Change)
       Slide 5 (stat overlap, coverage callout)
       Slide 6 (replace stats with before/after bar chart + coverage panel)
Outputs: EN_v3.pptx and ES_v3.pptx
"""

import sys, io
if __name__ == '__main__':
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches as I, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---------------------------------------------------------------------------
# BRAND PALETTE
# ---------------------------------------------------------------------------
COPPER      = RGBColor(0xB8, 0x73, 0x33)
COPPER_DARK = RGBColor(0x8B, 0x5A, 0x2B)
COPPER_PALE = RGBColor(0xF0, 0xE4, 0xD0)
CREAM_LT    = RGBColor(0xFA, 0xF6, 0xF1)
CREAM_MD    = RGBColor(0xF2, 0xED, 0xE6)
DARK        = RGBColor(0x2A, 0x1F, 0x10)
CHARCOAL    = RGBColor(0x2C, 0x2C, 0x2C)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GREY        = RGBColor(0x9E, 0x97, 0x8E)
GREY_LT     = RGBColor(0xCC, 0xC5, 0xBA)
PANEL_BG    = RGBColor(0x3A, 0x2A, 0x10)

# ---------------------------------------------------------------------------
# HELPERS
# ---------------------------------------------------------------------------
def set_bg(slide, rgb):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_tb(slide, text, l, t, w, h,
           size=10, bold=False, color=DARK,
           align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf = txb.text_frame
    tf.word_wrap = wrap
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_tb_multiline(slide, lines, l, t, w, h, wrap=True):
    """lines = list of (text, size, bold, color, align, italic)"""
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf = txb.text_frame
    tf.word_wrap = wrap
    first = True
    for (text, size, bold, color, align, italic) in lines:
        if first:
            para = tf.paragraphs[0]
            first = False
        else:
            para = tf.add_paragraph()
        para.alignment = align
        run = para.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return txb


def add_rect(slide, l, t, w, h, fill=None, line_color=None, lw=Pt(0.75)):
    shp = slide.shapes.add_shape(1, l, t, w, h)  # 1 = RECTANGLE
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line_color:
        shp.line.color.rgb = line_color
        shp.line.width = lw
    else:
        shp.line.fill.background()
    return shp


def add_circle(slide, cx, cy, d, fill=COPPER):
    l = cx - d // 2
    t = cy - d // 2
    shp = slide.shapes.add_shape(9, l, t, d, d)  # 9 = OVAL
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    return shp


def remove_shapes(slide, names):
    names_set = set(names)
    to_del = [s._element for s in slide.shapes if s.name in names_set]
    for el in to_del:
        el.getparent().remove(el)


def fix_text_colors(slide, shape_name, color):
    for shp in slide.shapes:
        if shp.name == shape_name and shp.has_text_frame:
            for para in shp.text_frame.paragraphs:
                for run in para.runs:
                    run.font.color.rgb = color


# ---------------------------------------------------------------------------
# BILINGUAL CONTENT
# ---------------------------------------------------------------------------

TIMELINE_EVENTS = {
    'en': [
        ("Jan 3, 2026",
         "Oil sector reopens to private investment",
         None),
        ("Jan 6, 2026",
         "Rystad Energy: $183B oil capex 2026-2040",
         "$41B in fabrication & construction alone"),
        ("Jan 9-10, 2026",
         "Repsol: triple production to 135,000 bpd",
         "New Orinoco Belt blocks = new infrastructure"),
        ("Jan 30, 2026",
         "Chevron Q4: 240,000 bpd, +50% in 18 months",
         "Expansion of existing facilities"),
        ("Feb 4, 2025",
         "Movistar/Telefonica: $500M investment",
         "805 5G nodes + 500 base stations = galvanized towers"),
        ("Mar 11, 2026",
         "Reuters: Chevron & Shell close first major deals",
         "Production contracts = immediate structural steel demand"),
        ("Sep 2025",
         "CONATEL: National Telecom Plan 2025-2031",
         "1,500 4G + 500 5G stations + 36,000 km fiber optic"),
        ("Mar 2026",
         "Venezuela reaches 1.1M bpd production",
         "Oil recovery underway — demand accelerates"),
    ],
    'es': [
        ("3 ene 2026",
         "El sector petrolero reabre a inversion privada",
         None),
        ("6 ene 2026",
         "Rystad Energy: US$183MM capex petrolero 2026-2040",
         "US$41MM en fabricacion y construccion"),
        ("9-10 ene 2026",
         "Repsol: triplicar produccion a 135.000 bpd",
         "Nuevos bloques en la Faja del Orinoco"),
        ("30 ene 2026",
         "Chevron Q4: 240.000 bpd, +50% en 18 meses",
         "Expansion de instalaciones existentes"),
        ("4 feb 2025",
         "Movistar/Telefonica: inversion de $500M",
         "805 nodos 5G + 500 estaciones = torres galvanizadas"),
        ("11 mar 2026",
         "Reuters: Chevron y Shell cierran primeros contratos",
         "Contratos de produccion = demanda inmediata de acero"),
        ("sep 2025",
         "CONATEL: Plan Nacional de Telecom 2025-2031",
         "1.500 4G + 500 5G + 36.000 km de fibra optica"),
        ("mar 2026",
         "Venezuela alcanza 1,1MM bpd de produccion",
         "Recuperacion petrolera en marcha"),
    ]
}

TIMELINE_FOOTER = {
    'en': "These are not projections — these are signed contracts and public announcements with verified dates.",
    'es': "Estas no son proyecciones — son contratos firmados y anuncios publicos con fechas verificadas.",
}

SLIDE5_COVERAGE = {
    'en': "PYGLARA's kettles handle every product category on this slide. At 300 TM/month, we need to capture less than 3% of the oil sector's annual galvanizing demand to operate at 100% of proven capacity. Only 2 active galvanizers serve all of Venezuela.",
    'es': "Las cubas de PYGLARA cubren todas las categorias de este slide. Con 300 TM/mes, necesitamos captar menos del 3% de la demanda anual de galvanizado del sector petrolero para operar al 100% de la capacidad comprobada. Solo 2 galvanizadores activos en toda Venezuela.",
}

SLIDE6 = {
    'en': {
        'chart_title': "",
        'bar1_label': "2024\nExisting",
        'bar2_label': "2031\nPlan",
        'bar1_val': 3000,
        'bar2_val': 5000,
        'y_unit': "Stations",
        'bar_annotation': "+2,000 new stations",
        'panel_header': "PYGLARA SERVES BOTH DEMAND STREAMS",
        'bullets': [
            "Tower steel demand: ~10,000 tons (2025-2031)",
            "Ground rod demand: ~18,000 copper rods nationwide",
            "At 300 TM/month: all tower steel covered in 3 months",
            "At 936 rods/day: all ground rods produced in 20 days",
            "Only company in Venezuela providing both: galvanized structure + copper grounding",
        ],
        'footer': "Every cell tower requires galvanized structural steel. Every base station requires copper-clad ground rods. PYGLARA is the only company in Venezuela that can supply both.",
    },
    'es': {
        'chart_title': "",
        'bar1_label': "2024\nExistentes",
        'bar2_label': "2031\nPlan",
        'bar1_val': 3000,
        'bar2_val': 5000,
        'y_unit': "Estaciones",
        'bar_annotation': "+2.000 nuevas estaciones",
        'panel_header': "PYGLARA ATIENDE AMBAS DEMANDAS",
        'bullets': [
            "Demanda de acero estructural: ~10.000 TM (2025-2031)",
            "Demanda de varillas de cobre: ~18.000 unidades",
            "Con 300 TM/mes: toda la demanda de acero cubierta en 3 meses",
            "Con 936 varillas/dia: toda la demanda en menos de 20 dias",
            "Unica empresa en Venezuela que provee ambos: acero + cobre",
        ],
        'footer': "Cada torre requiere estructura de acero galvanizado. Cada estacion requiere varillas de tierra cobreadas. PYGLARA es la unica empresa en Venezuela que puede proveer ambas.",
    }
}


# ---------------------------------------------------------------------------
# SLIDE 4 — Color revert + Vertical Timeline
# ---------------------------------------------------------------------------
def fix_slide4(slide, lang):
    # 1. Revert background to warm cream
    set_bg(slide, CREAM_MD)

    # 2. Fix title + subtitle text colors (were white on dark bg)
    fix_text_colors(slide, 'TextBox 2', DARK)
    fix_text_colors(slide, 'TextBox 3', COPPER)

    # 3. Remove all event TextBoxes (4 through 28)
    remove_shapes(slide, [f"TextBox {n}" for n in range(4, 29)])

    # ---- Layout constants (slide: 10" x 5.625") ----
    MARGIN_L   = I(0.6)
    LINE_X     = I(2.05)        # X of vertical timeline line (center)
    LINE_W     = Emu(32760)     # ~0.036" thin bar
    NODE_D     = I(0.14)        # Circle node diameter
    DATE_W     = I(1.35)        # Date column width (left of line)
    EVENT_L    = LINE_X + I(0.18)  # Left edge of event text
    EVENT_W    = I(10.0) - EVENT_L - I(0.25)  # Right margin 0.25"

    events  = TIMELINE_EVENTS[lang]
    N       = len(events)
    T_START = Emu(1_080_000)    # 1.18" — just below subtitle
    T_END   = Emu(4_250_000)    # 4.64" — above footer zone
    SPAN    = T_END - T_START
    STEP    = SPAN // N

    # 4. Draw vertical copper line
    line_top = T_START + STEP // 2
    line_h   = T_END - T_START - STEP // 2
    add_rect(slide,
             LINE_X - LINE_W // 2, line_top,
             LINE_W, line_h,
             fill=COPPER)

    # 5. Draw each event
    for idx, (date, title, impact) in enumerate(events):
        cy = T_START + idx * STEP + STEP // 2   # center Y of this row

        # Copper node circle
        add_circle(slide, LINE_X, cy, NODE_D)

        # Date label (right-aligned, left of line)
        add_tb(slide, date,
               MARGIN_L, cy - I(0.115),
               DATE_W - I(0.1), I(0.23),
               size=7.5, color=COPPER_DARK, align=PP_ALIGN.RIGHT, italic=False)

        if impact:
            # Title on top half, impact on bottom half
            add_tb(slide, title,
                   EVENT_L, cy - I(0.2),
                   EVENT_W, I(0.2),
                   size=9, bold=True, color=DARK)
            add_tb(slide, impact,
                   EVENT_L, cy + I(0.01),
                   EVENT_W, I(0.17),
                   size=8, color=COPPER_DARK, italic=True)
        else:
            # Single-line entry — center it vertically
            add_tb(slide, title,
                   EVENT_L, cy - I(0.13),
                   EVENT_W, I(0.27),
                   size=9.5, bold=True, color=COPPER)

    # 6. Footer
    add_tb(slide, TIMELINE_FOOTER[lang],
           MARGIN_L, I(4.7),
           I(8.8), I(0.4),
           size=8, italic=True, color=GREY)


# ---------------------------------------------------------------------------
# SLIDE 5 — Fix stat overlap + Coverage callout
# ---------------------------------------------------------------------------
def fix_slide5(slide, lang):
    # Fix overlapping right-panel stat boxes
    # TB30 ($10B+ materials): top=2,514,600, h=548,640 → bottom=3,063,240
    # TB31 (Materials label): was 3,017,520 → move to 3,063,240 (clean stack)
    # TB32 ($10B+ electrical): was 3,246,120 → move to 3,520,440 (bottom of TB31+h)
    # TB33 (Electrical label): was 3,749,040 → move to 4,069,080 (bottom of TB32+h)
    repositions = {
        'TextBox 31': 3_063_240,
        'TextBox 32': 3_520_440,
        'TextBox 33': 4_069_080,
    }
    for shp in slide.shapes:
        if shp.name in repositions:
            shp.top = repositions[shp.name]

    # Center-align all right-panel stat text boxes
    center_names = {'TextBox 28', 'TextBox 29', 'TextBox 30',
                    'TextBox 31', 'TextBox 32', 'TextBox 33'}
    for shp in slide.shapes:
        if shp.name in center_names and shp.has_text_frame:
            for para in shp.text_frame.paragraphs:
                para.alignment = PP_ALIGN.CENTER

    # Coverage callout strip at bottom
    # Content ends at ~4.59" after repositions. Slide height = 5.625"
    CALL_T = I(4.72)
    CALL_H = I(0.6)
    add_rect(slide, I(0.55), CALL_T, I(8.9), CALL_H,
             fill=COPPER_PALE, line_color=COPPER, lw=Pt(1))
    add_tb(slide, SLIDE5_COVERAGE[lang],
           I(0.7), CALL_T + I(0.04),
           I(8.6), CALL_H - I(0.08),
           size=8, color=DARK, align=PP_ALIGN.CENTER, italic=True)


# ---------------------------------------------------------------------------
# SLIDE 6 — Before/After Bar Chart + Coverage Panel
# ---------------------------------------------------------------------------
def fix_slide6(slide, lang):
    d = SLIDE6[lang]

    # 1. Remove existing stat boxes and separator
    remove_shapes(slide, [f"TextBox {n}" for n in range(3, 9)])
    remove_shapes(slide, ['Rectangle 10', 'TextBox 11', 'TextBox 12', 'TextBox 13'])

    # 2. Move TextBox 9 (the "both needs" paragraph) to footer position
    for shp in slide.shapes:
        if shp.name == 'TextBox 9' and shp.has_text_frame:
            shp.top    = I(5.05)
            shp.left   = I(0.6)
            shp.width  = I(8.8)
            shp.height = I(0.5)
            for para in shp.text_frame.paragraphs:
                para.alignment = PP_ALIGN.CENTER
                for run in para.runs:
                    run.font.size = Pt(8.5)
                    run.font.italic = True
                    run.font.color.rgb = GREY_LT

    # ---- Bar Chart (left panel) ----
    # Area: L=0.4", T=0.85", W=4.7", H=3.55"
    CL = I(0.4)
    CT = I(0.85)
    CW = I(4.7)
    CH = I(3.55)
    MAX_V = 6000
    BASE_Y = CT + CH          # bottom of chart (Y baseline)

    def bar_h(v):
        return Emu(int(CH * v / MAX_V))

    def bar_t(v):
        return BASE_Y - bar_h(v)

    # Chart area width capped before the right panel (panel starts at I(5.3))
    CHART_W = I(4.45)  # X-axis ends at I(0.75)+I(4.45)=I(5.20), 0.10" before panel

    # Axis lines
    add_rect(slide, CL + I(0.35), CT, Emu(18288), CH, fill=GREY)            # Y-axis
    add_rect(slide, CL + I(0.35), BASE_Y - Emu(18288), CHART_W, Emu(18288), fill=GREY)  # X-axis

    # Y-axis grid lines and tick labels
    for tick in [1000, 2000, 3000, 4000, 5000]:
        ty = bar_t(tick)
        add_rect(slide, CL + I(0.35), ty, CHART_W - I(0.1), Emu(9144), fill=RGBColor(0x55, 0x4A, 0x3F))  # grid line
        add_tb(slide, f"{tick // 1000}K",
               CL, ty - I(0.1), I(0.32), I(0.2),
               size=7, color=GREY_LT, align=PP_ALIGN.RIGHT)

    # Y-axis unit label moved above chart, below chart title to avoid overlap
    add_tb(slide, d['y_unit'],
           CL + I(0.38), CT - I(0.22), I(2.0), I(0.20),
           size=7.5, color=GREY_LT, italic=True)

    # Bar specs
    BAR_W = I(1.25)
    B1_X  = CL + I(0.65)   # 2024 bar
    B2_X  = CL + I(2.55)   # 2031 bar

    # 2024 bar (grey)
    b1_h = bar_h(d['bar1_val'])
    b1_t = bar_t(d['bar1_val'])
    add_rect(slide, B1_X, b1_t, BAR_W, b1_h, fill=GREY)
    add_tb(slide, f"{d['bar1_val']:,}",
           B1_X, b1_t - I(0.25), BAR_W, I(0.22),
           size=13, bold=True, color=GREY_LT, align=PP_ALIGN.CENTER)
    add_tb(slide, d['bar1_label'],
           B1_X, BASE_Y + I(0.07), BAR_W, I(0.4),
           size=8.5, color=GREY_LT, align=PP_ALIGN.CENTER)

    # 2031 bar (copper)
    b2_h = bar_h(d['bar2_val'])
    b2_t = bar_t(d['bar2_val'])
    add_rect(slide, B2_X, b2_t, BAR_W, b2_h, fill=COPPER)
    add_tb(slide, f"{d['bar2_val']:,}",
           B2_X, b2_t - I(0.25), BAR_W, I(0.22),
           size=13, bold=True, color=COPPER, align=PP_ALIGN.CENTER)
    add_tb(slide, d['bar2_label'],
           B2_X, BASE_Y + I(0.07), BAR_W, I(0.4),
           size=8.5, color=WHITE, align=PP_ALIGN.CENTER)

    # Delta annotation between bars
    ANNO_X = (B1_X + BAR_W + B2_X) // 2
    add_tb(slide, d['bar_annotation'],
           B1_X + BAR_W + I(0.05), b2_t - I(0.02),
           B2_X - (B1_X + BAR_W) - I(0.1), I(0.4),
           size=8, bold=True, color=COPPER, align=PP_ALIGN.CENTER)

    # Chart title removed — slide title already describes the topic

    # ---- Coverage Panel (right side) ----
    PL = I(5.3)
    PT = I(0.8)
    PW = I(4.35)
    PH = I(3.65)

    add_rect(slide, PL, PT, PW, PH, fill=PANEL_BG, line_color=COPPER, lw=Pt(1.5))

    # Panel header
    add_tb(slide, d['panel_header'],
           PL + I(0.15), PT + I(0.1), PW - I(0.3), I(0.35),
           size=9.5, bold=True, color=COPPER, align=PP_ALIGN.LEFT)

    # Divider line under header
    add_rect(slide,
             PL + I(0.15), PT + I(0.5), PW - I(0.3), Emu(13716),
             fill=COPPER_DARK)

    # Bullet items
    BULLET_TOP = PT + I(0.58)
    BULLET_STEP = I(0.57)
    BULLET_D = I(0.08)

    for i, bullet in enumerate(d['bullets']):
        bt = BULLET_TOP + i * BULLET_STEP
        # Copper dot
        add_circle(slide, PL + I(0.26), bt + I(0.09), BULLET_D, fill=COPPER)
        # Bullet text
        add_tb(slide, bullet,
               PL + I(0.38), bt,
               PW - I(0.53), BULLET_STEP - I(0.05),
               size=8.5, color=WHITE, wrap=True)


# ---------------------------------------------------------------------------
# SLIDE-BY-SLIDE AUDIT (for both versions)
# ---------------------------------------------------------------------------
def print_audit():
    print("\n" + "="*70)
    print("NARRATIVE AUDIT — Problem > Solution > Coverage Framework")
    print("="*70)
    audit = [
        ("1",  "Cover",                   "N/A",       "OK"),
        ("2",  "The Opportunity",         "Implicit",  "OK — summary stats present"),
        ("3",  "Why Now",                 "Demand",    "OK — add coverage note manually"),
        ("4",  "Demand Is Already Here",  "FIXED",     "Timeline rebuilt, Regime Change removed"),
        ("5",  "What Oil Sector Galvanizes","FIXED",   "Stat overlap fixed, coverage callout added"),
        ("6",  "Telecommunications",      "REBUILT",   "Bar chart + coverage panel added"),
        ("7",  "Why Galvanizing",         "Education", "OK — no coverage needed"),
        ("8",  "The Plant",               "Asset",     "OK"),
        ("9",  "Strategic Location",      "Logistics", "OK"),
        ("10", "Geographic Advantages",   "Logistics", "OK"),
        ("11", "World-Class Equipment",   "Asset",     "Could add: 7m kettle = 300 TM/month capacity"),
        ("12", "Pilling Advantage",       "Asset",     "OK"),
        ("13", "Products & Applications", "Product",   "RECOMMEND: add TM demand per category"),
        ("14", "Copper Electroplating",   "Product",   "OK — zero competitor angle strong"),
        ("15", "Quality & Certifications","Credential","OK"),
        ("16", "The Team",                "Team",      "OK"),
        ("17", "Track Record",            "Historical","OK — SENIAT-confirmed data strong"),
        ("18", "Current Inventory",       "Asset",     "OK"),
        ("19", "Competitive Position",    "Market",    "RECOMMEND: add market share math here"),
        ("20", "Galvanizing Economics",   "Financial", "OK — cost structure clear"),
        ("21", "Financial Projections",   "Financial", "OK — 3 scenarios"),
        ("22", "Break-Even",              "Financial", "OK — 10-20 TM/month is compelling"),
        ("23", "Capital Deployment",      "Financial", "OK — phased investment clear"),
        ("24", "Asset Valuation",         "Valuation", "OK"),
        ("25", "What's Included",         "Deal",      "OK"),
        ("26", "Next Steps",              "CTA",       "OK"),
        ("27", "Back Cover",              "N/A",       "OK"),
        ("28", "OFAC Compliance",         "Legal",     "OK"),
        ("29", "USD Repatriation",        "Legal",     "OK"),
        ("30", "Key Risks",               "Risk",      "OK"),
    ]
    for (num, name, type_, status) in audit:
        flag = ">> ACTION" if "RECOMMEND" in status or "FIXED" in status or "REBUILT" in status else "   "
        print(f"  {flag} Slide {num:>2}: {name:<30} [{type_:<12}] {status}")
    print()


# ---------------------------------------------------------------------------
# MAIN
# ---------------------------------------------------------------------------
if __name__ == '__main__':
    DECKS = [
        (
            'c:/Users/amont/Desktop/PYG/investor-ready/01-Pitch-Deck/PYGLARA_Pitch_Deck_EN.pptx',
            'c:/Users/amont/Desktop/PYG/investor-ready/01-Pitch-Deck/PYGLARA_Pitch_Deck_EN_v3.pptx',
            'en',
        ),
        (
            'c:/Users/amont/Desktop/PYG/investor-ready/01-Pitch-Deck/PYGLARA_Pitch_Deck.pptx',
            'c:/Users/amont/Desktop/PYG/investor-ready/01-Pitch-Deck/PYGLARA_Pitch_Deck_ES_v3.pptx',
            'es',
        ),
    ]

    for src, dst, lang in DECKS:
        print(f"\nProcessing {lang.upper()} deck: {src}")
        prs = Presentation(src)
        print(f"  Dimensions: {prs.slide_width/914400:.2f}\" x {prs.slide_height/914400:.2f}\"  ({len(prs.slides)} slides)")

        fix_slide4(prs.slides[3], lang)
        print(f"  [OK] Slide 4: background reverted, Regime Change removed, vertical timeline built")

        fix_slide5(prs.slides[4], lang)
        print(f"  [OK] Slide 5: stat overlaps fixed, text centered, coverage callout added")

        fix_slide6(prs.slides[5], lang)
        print(f"  [OK] Slide 6: stat boxes replaced with before/after bar chart + PYGLARA coverage panel")

        prs.save(dst)
        print(f"  [SAVED] {dst}")

    print_audit()
    print("\nDone. Both decks saved as v3.")
