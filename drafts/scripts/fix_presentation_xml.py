#!/usr/bin/env python3
"""
Fix presentation.xml by restoring it from the original PPTX and
applying string-level rId substitution. Then re-export to PDF.
"""
import sys, io, os, zipfile, re
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

BASE = r'c:\Users\amont\Desktop\PYG\investor-ready\01-Pitch-Deck'

PAIRS = [
    ('PYGLARA_Pitch_Deck_EN.pptx',  'PYGLARA_Pitch_Deck_EN_v3.pptx'),
    ('PYGLARA_Pitch_Deck.pptx',     'PYGLARA_Pitch_Deck_ES_v3.pptx'),
]

# rId remapping from our fix_pptx_rels2.py pass
RID_MAP = {
    'rId34': 'rId5',
    'rId35': 'rId6',
    'rId36': 'rId7',
    'rId5':  'rId8',
    'rId6':  'rId9',
    'rId7':  'rId10',
    'rId8':  'rId11',
    'rId9':  'rId12',
    'rId10': 'rId13',
    'rId11': 'rId14',
    'rId12': 'rId15',
    'rId13': 'rId16',
    'rId14': 'rId17',
    'rId15': 'rId18',
    'rId16': 'rId19',
    'rId17': 'rId20',
    'rId18': 'rId21',
    'rId19': 'rId22',
    'rId20': 'rId23',
    'rId21': 'rId24',
    'rId22': 'rId25',
    'rId23': 'rId26',
    'rId24': 'rId27',
    'rId25': 'rId28',
    'rId26': 'rId29',
    'rId27': 'rId30',
    'rId28': 'rId31',
}

def apply_rid_map(xml_str, rid_map):
    """Apply rId substitutions using word-boundary matching."""
    # Sort by length descending to avoid partial matches (rId10 before rId1)
    for old, new in sorted(rid_map.items(), key=lambda x: -len(x[0])):
        # Match only when surrounded by quotes (XML attribute values)
        xml_str = xml_str.replace('"' + old + '"', '"' + new + '"')
    return xml_str


for orig_name, v3_name in PAIRS:
    print(f'\nFixing {v3_name}')
    orig_path = os.path.join(BASE, orig_name)
    v3_path   = os.path.join(BASE, v3_name)

    # Get original presentation.xml (clean namespace structure)
    with open(orig_path, 'rb') as fh:
        orig_raw = fh.read()
    with zipfile.ZipFile(io.BytesIO(orig_raw)) as z:
        orig_prs_xml = z.read('ppt/presentation.xml').decode('utf-8')

    # Apply rId substitution in presentation.xml
    fixed_prs = apply_rid_map(orig_prs_xml, RID_MAP)

    # Verify
    old_rids_present = [r for r in ['rId34', 'rId35', 'rId36'] if r in fixed_prs]
    if old_rids_present:
        print(f'  ERROR: Old rIds still present: {old_rids_present}')
    else:
        print(f'  rId substitution OK — no old rIds remaining')

    # Sample check
    if '"rId5"' in fixed_prs:
        print(f'  rId5 present in sldIdLst (slide4 position)')
    if '"rId6"' in fixed_prs:
        print(f'  rId6 present (slide5 position)')

    # Load v3 zip and replace presentation.xml
    with open(v3_path, 'rb') as fh:
        v3_raw = fh.read()
    with zipfile.ZipFile(io.BytesIO(v3_raw)) as z:
        names = z.namelist()
        files = {n: z.read(n) for n in names}

    files['ppt/presentation.xml'] = fixed_prs.encode('utf-8')

    # Write back
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, 'w', zipfile.ZIP_DEFLATED) as zout:
        for name in names:
            zout.writestr(name, files[name])

    with open(v3_path, 'wb') as fh:
        fh.write(buf.getvalue())

    size_kb = os.path.getsize(v3_path) // 1024
    print(f'  Written: {size_kb} KB')

    # Verify with python-pptx
    from pptx import Presentation
    try:
        p = Presentation(v3_path)
        slides = len(p.slides)
        s4 = [t.strip() for s in p.slides[3].shapes if s.has_text_frame
              for t in [pr.text for pr in s.text_frame.paragraphs] if t.strip()]
        print(f'  python-pptx: {slides} slides, slide4 first text: {s4[0] if s4 else "empty"}')
    except Exception as e:
        print(f'  python-pptx error: {e}')

print('\n--- All XML fixes applied ---')
print('Now attempting PDF export via PowerPoint...')

import subprocess, time, win32com.client

for orig_name, v3_name in PAIRS:
    v3_path  = os.path.join(BASE, v3_name)
    pdf_name = v3_name.replace('.pptx', '.pdf')
    pdf_path = os.path.join(BASE, pdf_name)

    if os.path.exists(pdf_path):
        os.remove(pdf_path)

    proc = subprocess.Popen(
        [r'C:\Program Files\Microsoft Office\root\Office16\POWERPNT.EXE', v3_path]
    )
    print(f'\nLaunched PPT PID={proc.pid} for {v3_name}')

    ppt = None
    for attempt in range(12):
        time.sleep(2)
        try:
            ppt = win32com.client.GetActiveObject('PowerPoint.Application')
            count = ppt.Presentations.Count
            if count > 0:
                print(f'  Got PPT ({count} pres) on attempt {attempt+1}')
                break
        except:
            pass
    else:
        print('  Could not get PPT instance with open presentations')
        proc.kill()
        continue

    try:
        prs = ppt.Presentations(1)
        name = prs.Name
        slides = prs.Slides.Count
        print(f'  Pres: {name} — {slides} slides')
        if '[Repaired]' in name:
            print('  WARNING: Still in Repaired mode')
        prs.SaveAs(pdf_path, 32)
        size_kb = os.path.getsize(pdf_path) // 1024
        print(f'  [OK] {pdf_name} — {size_kb} KB')
        prs.Close()
    except Exception as e:
        print(f'  Export error: {e}')

    try:
        ppt.Quit()
    except:
        pass

    time.sleep(2)

print('\nDone.')
