#!/usr/bin/env python3
"""
Fix the ns0: namespace prefix issue introduced by ET.tostring() in the rels file.
Replace ns0:Relationships/ns0:Relationship with proper unnamespaced elements.
Also fix presentation.xml if it has ns0: prefixes.
"""
import sys, io, os, zipfile, re
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

BASE = r'c:\Users\amont\Desktop\PYG\investor-ready\01-Pitch-Deck'
DECKS = [
    'PYGLARA_Pitch_Deck_EN_v3.pptx',
    'PYGLARA_Pitch_Deck_ES_v3.pptx',
]

CORRECT_RELS_HEADER = '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'

def fix_ns(xml_str):
    """Remove ns0: prefix and fix xmlns declaration."""
    # Replace opening tag with proper one
    xml_str = re.sub(
        r'<ns0:Relationships[^>]+>',
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">',
        xml_str
    )
    xml_str = xml_str.replace('</ns0:Relationships>', '</Relationships>')
    xml_str = xml_str.replace('<ns0:Relationship ', '<Relationship ')
    xml_str = xml_str.replace('/>', ' />')  # normalize self-closing
    return xml_str

def fix_prs_xml(xml_str):
    """Fix presentation.xml if it has ns0: namespace prefixes introduced by ET."""
    # The sldIdLst update via ET.tostring may have changed namespace prefixes
    # Look for the problematic pattern
    if 'ns0:' not in xml_str and 'ns1:' not in xml_str:
        return xml_str  # No fix needed
    # This is more complex — let's just check
    print('  WARNING: presentation.xml has ns prefixes, may need manual fix')
    return xml_str


for deck in DECKS:
    path = os.path.join(BASE, deck)
    print(f'\n{deck}')

    with open(path, 'rb') as fh:
        raw = fh.read()

    with zipfile.ZipFile(io.BytesIO(raw), 'r') as zin:
        names = zin.namelist()
        files = {n: zin.read(n) for n in names}

    # Fix rels
    rels_str = files['ppt/_rels/presentation.xml.rels'].decode('utf-8')
    fixed_rels = fix_ns(rels_str)
    files['ppt/_rels/presentation.xml.rels'] = fixed_rels.encode('utf-8')

    # Show before/after
    print('  RELS before:', rels_str[38:90])
    print('  RELS after: ', fixed_rels[38:90])

    # Check presentation.xml
    prs_str = files['ppt/presentation.xml'].decode('utf-8')
    fixed_prs = fix_prs_xml(prs_str)
    files['ppt/presentation.xml'] = fixed_prs.encode('utf-8')

    # Write back
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, 'w', zipfile.ZIP_DEFLATED) as zout:
        for name in names:
            zout.writestr(name, files[name])

    with open(path, 'wb') as fh:
        fh.write(buf.getvalue())

    print(f'  Written: {os.path.getsize(path)//1024} KB')

    # Quick check with python-pptx
    from pptx import Presentation
    try:
        p = Presentation(path)
        print(f'  python-pptx: {len(p.slides)} slides')
        # Verify slide 4 content
        s4_texts = [t.strip() for s in p.slides[3].shapes if s.has_text_frame
                    for t in [p2.text for p2 in s.text_frame.paragraphs] if t.strip()]
        print(f'  Slide 4 first texts: {s4_texts[:3]}')
    except Exception as e:
        print(f'  Error: {e}')

print('\nDone. Now exporting PDFs...')

import subprocess, time, shutil

for deck in DECKS:
    pptx_path = os.path.join(BASE, deck)
    pdf_name  = deck.replace('.pptx', '.pdf')
    pdf_path  = os.path.join(BASE, pdf_name)

    # Remove old PDF if any
    if os.path.exists(pdf_path):
        os.remove(pdf_path)

    # Launch PPTX directly in PowerPoint
    proc = subprocess.Popen(
        [r'C:\Program Files\Microsoft Office\root\Office16\POWERPNT.EXE', pptx_path]
    )
    print(f'\nLaunched PPT for {deck} (PID {proc.pid}), waiting 8s...')
    time.sleep(8)

    # Get instance and export
    import win32com.client
    try:
        ppt = win32com.client.GetActiveObject('PowerPoint.Application')
        print(f'  PPT active, {ppt.Presentations.Count} pres')
        if ppt.Presentations.Count > 0:
            prs = ppt.Presentations(1)
            print(f'  Pres: {prs.Name} — {prs.Slides.Count} slides')
            prs.SaveAs(pdf_path, 32)
            prs.Close()
            print(f'  [OK] {pdf_name} — {os.path.getsize(pdf_path)//1024} KB')
        ppt.Quit()
    except Exception as e:
        print(f'  Error: {e}')
        try: ppt.Quit()
        except: pass

    time.sleep(2)
